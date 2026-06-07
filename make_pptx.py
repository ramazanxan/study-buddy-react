"""
StudyBuddy — генератор презентации PowerPoint
Цвета берутся с сайта: тёмный фон + зелёный акцент
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import urllib.request, io, os

# ── Цветовая палитра сайта ──────────────────────────────────────────
BG_DARK   = RGBColor(0x0D, 0x1A, 0x0F)   # очень тёмный зеленовато-чёрный
BG_CARD   = RGBColor(0x13, 0x24, 0x16)   # чуть светлее для карточек
GREEN     = RGBColor(0x22, 0xC5, 0x5E)   # основной зелёный с сайта (#22c55e)
GREEN_LT  = RGBColor(0x86, 0xEF, 0xAC)   # светло-зелёный
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x9C, 0xA3, 0xAF)
SLIDE_W   = Inches(13.33)
SLIDE_H   = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H


def blank_slide(prs):
    """Добавить пустой слайд (без плейсхолдеров)."""
    layout = prs.slide_layouts[6]   # blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color: RGBColor):
    """Залить фон слайда сплошным цветом."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color: RGBColor, alpha=None):
    """Добавить прямоугольник-заливку."""
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


def txb(slide, text, l, t, w, h,
        size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, word_wrap=True):
    """Добавить текстовый блок."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = word_wrap
    tf = box.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_num_badge(slide, num_str, l, t, size=0.55):
    """Номер слайда — зелёный кружок с цифрой."""
    circ = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(size), Inches(size))
    circ.fill.solid()
    circ.fill.fore_color.rgb = GREEN
    circ.line.fill.background()
    tf = circ.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num_str
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = WHITE


def accent_bar(slide, t=0.22, h=0.045):
    """Горизонтальная зелёная линия-акцент."""
    bar = slide.shapes.add_shape(1, Inches(0.6), Inches(t), Inches(1.2), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()


# ══════════════════════════════════════════════════════════════════════
# СЛАЙД 1 — ОБЛОЖКА / HERO
# ══════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
fill_bg(s, BG_DARK)

# декоративные зелёные прямоугольники (орбы)
orb1 = s.shapes.add_shape(1, Inches(9.5), Inches(-1.0), Inches(5.5), Inches(5.5))
orb1.fill.solid(); orb1.fill.fore_color.rgb = RGBColor(0x16, 0x3A, 0x1C)
orb1.line.fill.background()

orb2 = s.shapes.add_shape(1, Inches(-1.0), Inches(4.5), Inches(4.0), Inches(4.0))
orb2.fill.solid(); orb2.fill.fore_color.rgb = RGBColor(0x16, 0x3A, 0x1C)
orb2.line.fill.background()

# значок КГТУ (скачиваем и вставляем)
logo_url = 'https://enactus.kg/wp-content/uploads/2022/04/kgtu-logo.png'
try:
    logo_data = urllib.request.urlopen(logo_url, timeout=8).read()
    s.shapes.add_picture(io.BytesIO(logo_data), Inches(5.9), Inches(0.6), Inches(1.5), Inches(1.5))
except Exception:
    pass

# бейдж
badge = add_rect(s, 3.5, 2.35, 6.3, 0.42, RGBColor(0x16, 0x3A, 0x1C))
txb(s, '● Только для студентов КГТУ им. И. Раззакова', 3.55, 2.37,
    6.2, 0.38, size=11, color=GREEN_LT, align=PP_ALIGN.CENTER)

# логотип
txb(s, '✦  StudyBuddy', 4.5, 2.9, 4.3, 0.5,
    size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# заголовок hero
txb(s, 'Найди своих', 1.2, 3.45, 10.9, 1.1,
    size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, 'по учёбе', 1.2, 4.4, 10.9, 1.1,
    size=72, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# subtitle
txb(s, 'Платформа для студентов КГТУ — находи напарников, наставников и достигай целей вместе.',
    1.5, 5.55, 10.3, 0.7, size=16, color=GRAY, align=PP_ALIGN.CENTER)

# stats strip
stats = [('18 000+', 'Студентов КГТУ'), ('10', 'Институтов'),
         ('70+', 'Специальностей'), ('1954', 'Год основания')]
sx = 1.0
for num, lbl in stats:
    txb(s, num,  sx, 6.4, 2.7, 0.52, size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    txb(s, lbl,  sx, 6.85, 2.7, 0.38, size=11, color=GRAY, align=PP_ALIGN.CENTER)
    sx += 2.9
    if sx < 11:
        div = s.shapes.add_shape(1, Inches(sx-0.08), Inches(6.45), Inches(0.02), Inches(0.75))
        div.fill.solid(); div.fill.fore_color.rgb = RGBColor(0x2D, 0x4D, 0x32)
        div.line.fill.background()

# номер слайда
add_num_badge(s, '01', 0.25, 0.25)


# ══════════════════════════════════════════════════════════════════════
# СЛАЙД 2 — ПРОБЛЕМА
# ══════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
fill_bg(s, BG_DARK)

add_num_badge(s, '02', 0.25, 0.25)
accent_bar(s)

txb(s, 'ПРОБЛЕМА', 0.6, 0.35, 4, 0.35, size=11, color=GREEN, bold=True)
txb(s, 'В КГТУ тысячи\nстудентов.', 0.6, 0.75, 7, 1.8,
    size=58, bold=True, color=WHITE)
txb(s, 'Почему так сложно найти своих?', 0.6, 2.65, 7, 0.55,
    size=20, color=GRAY)

# правый декор — карточки "найти команду"
card = add_rect(s, 7.8, 1.0, 5.0, 5.2, BG_CARD)
txb(s, '🔍  Найти напарника по учёбе', 8.05, 1.25, 4.5, 0.5,
    size=15, color=WHITE, bold=True)
txb(s, '👥  Найти наставника (Агай/Эже)', 8.05, 1.95, 4.5, 0.5,
    size=15, color=WHITE, bold=True)
txb(s, '🎯  Достигать целей вместе', 8.05, 2.65, 4.5, 0.5,
    size=15, color=WHITE, bold=True)
txb(s, '🤝  Заключать пакты и держать слово', 8.05, 3.35, 4.5, 0.5,
    size=15, color=WHITE, bold=True)
txb(s, '🏆  Делиться победами', 8.05, 4.05, 4.5, 0.5,
    size=15, color=WHITE, bold=True)

txb(s, 'Мы это меняем.', 0.6, 3.5, 7, 0.6,
    size=26, bold=True, color=GREEN)


# ══════════════════════════════════════════════════════════════════════
# СЛАЙД 3 — ЧТО ТАКОЕ STUDYBUDDY
# ══════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
fill_bg(s, BG_DARK)

add_num_badge(s, '03', 0.25, 0.25)
accent_bar(s)

txb(s, 'РЕШЕНИЕ', 0.6, 0.35, 4, 0.35, size=11, color=GREEN, bold=True)
txb(s, 'Мы создали место,\nгде студенты\nнаходят друг друга.', 0.6, 0.75, 7.5, 2.4,
    size=48, bold=True, color=WHITE)

txb(s, 'Не соцсеть. Не мессенджер. Не маркетплейс.\nЭто платформа для тех, кто хочет создавать.',
    0.6, 3.3, 7.5, 0.85, size=18, color=GRAY)

# 3 иконки-фишки
features3 = [
    ('✦', 'Умный подбор', 'Искра — алгоритм подберёт\nнапарников по интересам\nи целям'),
    ('🤝', 'Пакты', 'Договор с напарником —\nдержи слово до дедлайна'),
    ('🧭', 'Наставники', 'Агай/Эже — старшекурсники\nдля младших'),
]
fx = 0.5
for icon, title, desc in features3:
    fc = add_rect(s, fx, 4.3, 3.9, 2.7, BG_CARD)
    txb(s, icon, fx+0.2, 4.45, 0.7, 0.6, size=28)
    txb(s, title, fx+0.9, 4.5, 2.8, 0.5, size=18, bold=True, color=GREEN)
    txb(s, desc, fx+0.2, 5.05, 3.6, 1.7, size=13, color=GRAY)
    fx += 4.1


# ══════════════════════════════════════════════════════════════════════
# СЛАЙД 4 — КАК ЭТО РАБОТАЕТ
# ══════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
fill_bg(s, BG_DARK)

add_num_badge(s, '04', 0.25, 0.25)
accent_bar(s)

txb(s, 'КАК ЭТО РАБОТАЕТ', 0.6, 0.35, 6, 0.35, size=11, color=GREEN, bold=True)
txb(s, 'Три простых шага к\nпродуктивной учёбе', 0.6, 0.75, 12, 1.6,
    size=50, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, 'Присоединяйся и начни прямо сейчас', 0.6, 2.45, 12.1, 0.45,
    size=16, color=GRAY, align=PP_ALIGN.CENTER)

steps = [
    ('1', '📝', 'Создай профиль', 'Расскажи о себе, факультете\nи своей цели на учёбу.'),
    ('2', '✦',  'Найди Искру',   'Алгоритм подберёт студентов\nс общими интересами и целями.'),
    ('3', '🚀', 'Достигайте вместе', 'Заключайте пакты, ставьте цели\nи празднуйте победы.'),
]
sx = 0.5
for n, icon, title, desc in steps:
    sc = add_rect(s, sx, 3.1, 3.9, 3.8, BG_CARD)
    # номер
    nb = s.shapes.add_shape(1, Inches(sx+0.2), Inches(3.25), Inches(0.55), Inches(0.55))
    nb.fill.solid(); nb.fill.fore_color.rgb = GREEN; nb.line.fill.background()
    tf = nb.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run(); r.text = n
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE
    txb(s, icon, sx+0.9, 3.25, 2.8, 0.55, size=26)
    txb(s, title, sx+0.2, 4.0, 3.5, 0.55, size=20, bold=True, color=WHITE)
    txb(s, desc, sx+0.2, 4.6, 3.5, 2.0, size=14, color=GRAY)
    sx += 4.15


# ══════════════════════════════════════════════════════════════════════
# СЛАЙД 5 — ВОЗМОЖНОСТИ (ВСЁ СРАЗУ)
# ══════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
fill_bg(s, BG_DARK)

add_num_badge(s, '05', 0.25, 0.25)
accent_bar(s)

txb(s, 'ВОЗМОЖНОСТИ', 0.6, 0.35, 5, 0.35, size=11, color=GREEN, bold=True)
txb(s, 'Всё сразу.', 0.6, 0.75, 7, 0.95, size=62, bold=True, color=WHITE)
txb(s, 'Всё для студентов КГТУ — в одном месте.', 0.6, 1.75, 8, 0.45,
    size=17, color=GRAY)

features_all = [
    ('✦', 'Искра',        'Умный подбор партнёров по\nинтересам, целям и факультету.',  GREEN),
    ('🤝', 'Пакт',         'Договор с напарником и\nконтроль дедлайнов.',               RGBColor(0x60, 0x80, 0xEB)),
    ('🧭', 'Агай / Эже',   'Наставничество: старшекурсники\nдля младших студентов.',     RGBColor(0x7C, 0x3A, 0xED)),
    ('🙋', 'Я тоже хочу', 'Присоединяйся к чужим целям\nи иди к ним вместе.',           RGBColor(0x0E, 0xA5, 0xE9)),
    ('🏆', 'Доска побед',  'Делитесь достижениями\nи вдохновляйте сообщество.',          RGBColor(0xD9, 0x77, 0x06)),
    ('🛍️', 'Маркетплейс', 'Покупайте и продавайте книги,\nтехнику и вещи студентам.',  GREEN),
]
col, row = 0, 0
for icon, name, desc, color in features_all:
    fx = 0.4 + col * 4.25
    fy = 2.35 + row * 2.45
    fc = add_rect(s, fx, fy, 4.0, 2.2, BG_CARD)
    # цветной бейдж-иконка
    ib = s.shapes.add_shape(1, Inches(fx+0.2), Inches(fy+0.2), Inches(0.62), Inches(0.62))
    ib.fill.solid(); ib.fill.fore_color.rgb = color; ib.line.fill.background()
    txb(s, icon, fx+0.22, fy+0.22, 0.58, 0.58, size=20, align=PP_ALIGN.CENTER)
    txb(s, name, fx+0.95, fy+0.27, 2.85, 0.5, size=16, bold=True, color=WHITE)
    txb(s, desc, fx+0.2, fy+0.9, 3.65, 1.1, size=12, color=GRAY)
    col += 1
    if col == 3:
        col = 0; row += 1


# ══════════════════════════════════════════════════════════════════════
# СЛАЙД 6 — МАРКЕТПЛЕЙС
# ══════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
fill_bg(s, BG_DARK)

add_num_badge(s, '06', 0.25, 0.25)
accent_bar(s)

txb(s, 'МАРКЕТПЛЕЙС', 0.6, 0.35, 5, 0.35, size=11, color=GREEN, bold=True)
txb(s, 'Маркетплейс\nдля студентов.', 0.6, 0.75, 6.5, 1.8, size=50, bold=True, color=WHITE)
txb(s, 'Покупай и продавай то,\nчто нужно для учёбы.', 0.6, 2.65, 5.5, 0.85,
    size=18, color=GRAY)

items_list = ['✓  Конспекты', '✓  Книги', '✓  Планшеты', '✓  Чертежи и другое']
iy = 3.65
for item in items_list:
    txb(s, item, 0.6, iy, 4, 0.45, size=16, color=GREEN_LT, bold=True)
    iy += 0.5

txb(s, '💰  Оплата в сомах (KGS)', 0.6, 5.95, 5, 0.45, size=15, color=GRAY)

# витрина товаров
showcase = [
    ('📄', 'Конспект по Алгоритмам', '200 сом'),
    ('📗', 'Учебник Clean Code',     '450 сом'),
    ('💻', 'Планшет HUION H640P',    '2 500 сом'),
    ('📐', 'Сборник задач по матану','150 сом'),
    ('🎧', 'Наушники Sony WH-1000XM4','1 500 сом'),
    ('📐', 'Чертежи AutoCAD',        '100 сом'),
]
gx, gy = 7.2, 0.9
col = 0
for icon, name, price in showcase:
    cx = gx + col * 2.05
    cy = gy + (0 if col < 3 else 2.3)
    if col == 3: col = 0; gx = 7.2
    cc = add_rect(s, cx, cy, 1.9, 2.0, BG_CARD)
    txb(s, icon, cx+0.7, cy+0.2, 0.7, 0.7, size=26, align=PP_ALIGN.CENTER)
    txb(s, name,  cx+0.1, cy+1.0, 1.7, 0.7, size=10, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, price, cx+0.1, cy+1.65, 1.7, 0.35, size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    col += 1


# ══════════════════════════════════════════════════════════════════════
# СЛАЙД 7 — НАСТАВНИЧЕСТВО (АГАЙ/ЕЖЕ)
# ══════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
fill_bg(s, BG_DARK)

add_num_badge(s, '07', 0.25, 0.25)
accent_bar(s)

txb(s, 'НАСТАВНИЧЕСТВО', 0.6, 0.35, 5, 0.35, size=11, color=GREEN, bold=True)
txb(s, 'Агай & Эже —\nнайди наставника.', 0.6, 0.75, 8, 1.8,
    size=52, bold=True, color=WHITE)
txb(s, 'Старшекурсники становятся наставниками для младших.\nПомощь в учёбе, карьере и проектах.',
    0.6, 2.7, 7, 0.75, size=17, color=GRAY)

mentor_cards = [
    ('🧭', 'Академическая помощь', 'Помощь с предметами,\nпересдачами и курсовыми.'),
    ('💼', 'Карьера и стажировки', 'Советы по поиску работы\nи прохождению практики.'),
    ('🚀', 'Совместные проекты',   'Старт стартапов и\nхакатонов с наставником.'),
]
mx = 0.5
for icon, title, desc in mentor_cards:
    mc = add_rect(s, mx, 3.65, 3.9, 3.2, BG_CARD)
    txb(s, icon, mx+0.2, 3.85, 0.7, 0.7, size=30)
    txb(s, title, mx+0.2, 4.65, 3.5, 0.55, size=18, bold=True, color=WHITE)
    txb(s, desc, mx+0.2, 5.25, 3.5, 1.4, size=13, color=GRAY)
    mx += 4.15

txb(s, 'Стань наставником →', 9.2, 0.8, 4.0, 0.5,
    size=17, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════
# СЛАЙД 8 — ИСТОРИИ УСПЕХА
# ══════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
fill_bg(s, BG_DARK)

add_num_badge(s, '08', 0.25, 0.25)
accent_bar(s)

txb(s, 'ИСТОРИИ УСПЕХА', 0.6, 0.35, 6, 0.35, size=11, color=GREEN, bold=True)
txb(s, 'Найди команду\nдля своей идеи.', 0.6, 0.75, 8, 1.8,
    size=52, bold=True, color=WHITE)
txb(s, 'Реальные победы наших студентов', 0.6, 2.65, 8, 0.45,
    size=18, color=GRAY)

wins = [
    ('🚀', 'Данияр & Чолпон',
     'Запустили MVP стартапа и получили первых 50 пользователей!'),
    ('🥈', 'Бекзат',
     'Занял 2 место на университетском хакатоне'),
    ('📱', 'Айзада',
     'Выпустила первое мобильное приложение в сторе'),
]
wx = 0.5
for emoji, name, text in wins:
    wc = add_rect(s, wx, 3.35, 3.9, 3.65, BG_CARD)
    txb(s, emoji, wx+0.2, 3.5, 0.8, 0.8, size=36)
    txb(s, text, wx+0.2, 4.45, 3.5, 1.5, size=15, color=WHITE)
    txb(s, '— ' + name, wx+0.2, 5.9, 3.5, 0.45, size=13, color=GREEN, bold=True)
    wx += 4.15

txb(s, '"Твоя идея заслуживает\nправильных людей."', 8.5, 2.5, 4.5, 1.4,
    size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════
# СЛАЙД 9 — МАСШТАБ: КГТУ
# ══════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
fill_bg(s, BG_DARK)

add_num_badge(s, '09', 0.25, 0.25)
accent_bar(s)

txb(s, 'МАСШТАБ', 0.6, 0.35, 4, 0.35, size=11, color=GREEN, bold=True)
txb(s, '1 КГТУ', 0.6, 0.75, 8, 1.0, size=70, bold=True, color=WHITE)

txb(s, '18 000+', 0.6, 1.85, 6, 0.9, size=64, bold=True, color=GREEN)
txb(s, 'Студентов', 0.6, 2.7, 6, 0.5, size=22, color=WHITE)

# ∞ бесконечное количество идей
inf = s.shapes.add_shape(1, Inches(0.6), Inches(3.4), Inches(1.2), Inches(0.7))
inf.fill.solid(); inf.fill.fore_color.rgb = GREEN; inf.line.fill.background()
txb(s, '∞', 0.6, 3.4, 1.2, 0.7, size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, 'Бесконечное количество идей', 1.95, 3.5, 6, 0.55, size=20, color=WHITE, bold=True)

# правый столбец — статистика
stats2 = [('10', 'Институтов'), ('70+', 'Специальностей'), ('1954', 'Год основания')]
sy = 0.7
for num, lbl in stats2:
    sc = add_rect(s, 9.0, sy, 3.8, 1.65, BG_CARD)
    txb(s, num, 9.0, sy+0.1, 3.8, 0.9, size=44, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    txb(s, lbl, 9.0, sy+0.95, 3.8, 0.45, size=14, color=GRAY, align=PP_ALIGN.CENTER)
    sy += 1.8

txb(s, 'КГТУ им. И. Раззакова · г. Бишкек · Основан в 1954',
    0.6, 6.75, 12.1, 0.45, size=13, color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# СЛАЙД 10 — ПРИСОЕДИНЯЙСЯ / CTA
# ══════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
fill_bg(s, BG_DARK)

# большой зелёный акцент-фон
bg2 = s.shapes.add_shape(1, Inches(0), Inches(1.8), Inches(13.33), Inches(4.2))
bg2.fill.solid(); bg2.fill.fore_color.rgb = RGBColor(0x14, 0x32, 0x18)
bg2.line.fill.background()

add_num_badge(s, '10', 0.25, 0.25)

txb(s, 'ПРИСОЕДИНЯЙСЯ.', 0.6, 0.3, 12.1, 1.4,
    size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txb(s, 'Для студентов КГТУ. От студентов КГТУ.', 0.6, 2.1, 12.1, 0.6,
    size=22, color=GREEN_LT, align=PP_ALIGN.CENTER)

txb(s, 'Найди напарников · Заключи пакт · Достигни цели',
    0.6, 2.85, 12.1, 0.5, size=17, color=GRAY, align=PP_ALIGN.CENTER)

# кнопка
btn = add_rect(s, 4.4, 3.6, 4.5, 0.75, GREEN)
txb(s, 'Создать профиль →', 4.4, 3.65, 4.5, 0.65,
    size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# команда
txb(s, 'Разработано студентами ИИТ, КЛ-3-25', 0.6, 4.8, 12.1, 0.4,
    size=13, color=GRAY, align=PP_ALIGN.CENTER)

authors = 'Айдай Токтоналиева  ·  Ариетта Асанкадырова  ·  Арзыкан Ражапбаева  ·  Элеонора Орозалиева  ·  Ислам Шаршенбеков'
txb(s, authors, 0.6, 5.2, 12.1, 0.4, size=11, color=GRAY, align=PP_ALIGN.CENTER)

# КГТУ логотип внизу
try:
    s.shapes.add_picture(io.BytesIO(logo_data), Inches(5.9), Inches(5.8), Inches(1.5), Inches(1.5))
except Exception:
    txb(s, '⚙️  КГТУ', 5.5, 5.8, 2.3, 0.7, size=28, color=GREEN, align=PP_ALIGN.CENTER)

txb(s, '✦  StudyBuddy', 0.6, 6.95, 12.1, 0.45,
    size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── Сохранение ──────────────────────────────────────────────────────
out = r'C:\Users\user\study-buddy-react\StudyBuddy_Presentation.pptx'
prs.save(out)
print(f'OK: {out}')
